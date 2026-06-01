from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv
import uvicorn

load_dotenv()

FILES_BASE = Path(os.getenv("FILES_BASE", "C:/files"))
OUTPUT_BASE = Path(os.getenv("OUTPUT_BASE", "C:/output"))
PORT = int(os.getenv("PORT", "9000"))

app = FastAPI(title="PPTX Generator")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
TEXT_LEFT   = Inches(0.4)
TEXT_TOP    = Inches(0.6)
TEXT_W      = Inches(6.2)
TEXT_H      = Inches(6.3)
IMG_LEFT    = Inches(7.0)
IMG_TOP     = Inches(0.4)
IMG_W       = Inches(6.0)
IMG_H       = Inches(6.7)

BG_COLOR  = RGBColor(0x0D, 0x1B, 0x2A)   # koyu lacivert
TITLE_CLR = RGBColor(0xFF, 0xFF, 0xFF)
BODY_CLR  = RGBColor(0xCC, 0xDD, 0xEE)
HEAD_CLR  = RGBColor(0x4F, 0xC3, 0xF7)   # açık mavi başlık


class PPTXRequest(BaseModel):
    proje_adi: str


KONU_RE = re.compile(r"^Konu(\d+)$")
SLIDE_RE = re.compile(r"^Slide(\d+)\.txt$")


def list_konular(proje_dir: Path):
    """proje_dir altındaki Konu{N} klasörlerini N'e göre sıralı döndürür."""
    konular = []
    for child in proje_dir.iterdir():
        if not child.is_dir():
            continue
        m = KONU_RE.match(child.name)
        if m:
            konular.append((int(m.group(1)), child))
    konular.sort(key=lambda t: t[0])
    return konular


def list_slides(konu_dir: Path):
    """konu_dir altındaki Slide{M}.txt dosyalarını M'e göre sıralı döndürür."""
    slides = []
    for child in konu_dir.iterdir():
        if not child.is_file():
            continue
        m = SLIDE_RE.match(child.name)
        if m:
            slides.append((int(m.group(1)), child))
    slides.sort(key=lambda t: t[0])
    return slides


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text: str, slide_type: str):
    txBox = slide.shapes.add_textbox(TEXT_LEFT, TEXT_TOP, TEXT_W, TEXT_H)
    tf = txBox.text_frame
    tf.word_wrap = True

    if slide_type == "headline":
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = TITLE_CLR
        p.alignment = PP_ALIGN.LEFT
    else:
        lines = [l for l in text.split("\n") if l.strip()]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            is_header = not line.startswith("-") and i == 0
            run.text = line.lstrip("- ").strip()
            run.font.size = Pt(22) if is_header else Pt(17)
            run.font.bold = is_header
            run.font.color.rgb = HEAD_CLR if is_header else BODY_CLR
            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.LEFT


def add_slide(prs: Presentation, konu_dir: Path, slide_no: int, text: str, slide_type: str):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_COLOR)
    add_text_box(slide, text, slide_type)

    img_path = konu_dir / f"imageSlide{slide_no}.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), IMG_LEFT, IMG_TOP, IMG_W, IMG_H)

    return slide


@app.post("/generate-pptx")
def generate_pptx(req: PPTXRequest):
    proje_dir = FILES_BASE / req.proje_adi
    if not proje_dir.is_dir():
        raise HTTPException(status_code=404, detail=f"Proje klasörü bulunamadı: {req.proje_adi}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_count = 0
    for konu_id, konu_dir in list_konular(proje_dir):
        for slide_no, slide_file in list_slides(konu_dir):
            text = slide_file.read_text(encoding="utf-8")
            slide_type = "headline" if slide_no == 1 else "bullet"
            add_slide(prs, konu_dir, slide_no, text, slide_type)
            slide_count += 1

    if slide_count == 0:
        raise HTTPException(status_code=404, detail=f"Projede slide içeriği bulunamadı: {req.proje_adi}")

    output_dir = OUTPUT_BASE / req.proje_adi
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{req.proje_adi}.pptx"
    prs.save(str(output_path))

    return {"status": "ok", "path": str(output_path), "slide_count": slide_count}


@app.get("/health")
def health():
    return {"status": "ok"}


if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=PORT, reload=True)
